# Databricks notebook source
# MAGIC %md
# MAGIC # Build the demo PPTX (Databricks theme + speaker notes)
# MAGIC Generates NorthStar_Brand_Copilot_Demo.pptx into a UC Volume so it can be downloaded locally.

# COMMAND ----------
# MAGIC %pip install python-pptx --quiet
# COMMAND ----------
dbutils.library.restartPython()

# COMMAND ----------
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CAT, SCH = "REPLACE_WITH_CATALOG", "northstar_cpg"
spark.sql(f"CREATE VOLUME IF NOT EXISTS {CAT}.{SCH}.assets")
OUT_DIR = f"/Volumes/{CAT}/{SCH}/assets"
OUT = f"{OUT_DIR}/NorthStar_Brand_Copilot_Demo.pptx"

# --- Databricks brand palette ---
LAVA  = RGBColor(0xFF, 0x36, 0x21)   # Databricks red/orange
NAVY  = RGBColor(0x1B, 0x31, 0x39)   # oat dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x5A, 0x66, 0x6B)
LGRAY = RGBColor(0xF5, 0xF5, 0xF3)
GREEN = RGBColor(0x00, 0xA9, 0x72)
BLUE  = RGBColor(0x22, 0x72, 0xB4)
FONT  = "DM Sans"          # Databricks font; falls back gracefully

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set(tf_or_para, text, size, color, bold=False, font=FONT):
    p = tf_or_para
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def rect(slide, l, t, w, h, fill, line=None, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def box_text(shp, text, size, color, bold=True):
    tf = shp.text_frame; tf.word_wrap = True
    _set(tf.paragraphs[0], text, size, color, bold)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(title, kicker, bullets, speaker):
    """bullets: list of (text, level) ; level 0 = main, 1 = sub."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE, rounded=False)              # white bg
    rect(s, 0, 0, Inches(0.18), SH, LAVA, rounded=False)     # left lava bar
    # kicker + title
    k = textbox(s, Inches(0.7), Inches(0.45), Inches(11.8), Inches(0.4))
    _set(k.paragraphs[0], kicker.upper(), 13, LAVA, bold=True)
    tt = textbox(s, Inches(0.7), Inches(0.8), Inches(11.8), Inches(1.0))
    _set(tt.paragraphs[0], title, 30, NAVY, bold=True)
    rect(s, Inches(0.72), Inches(1.78), Inches(1.1), Inches(0.06), LAVA, rounded=False)
    # body
    body = textbox(s, Inches(0.72), Inches(2.1), Inches(11.9), Inches(4.9))
    first = True
    for text, lvl in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        bullet = ("•  " if lvl == 0 else "–  ") + text
        _set(p, bullet, 18 if lvl == 0 else 15, NAVY if lvl == 0 else GRAY,
             bold=(lvl == 0))
        p.level = lvl; p.space_after = Pt(8)
        if lvl == 1:
            p.runs[0].font.bold = False
    # footer
    f = textbox(s, Inches(0.7), Inches(7.0), Inches(11.9), Inches(0.4))
    _set(f.paragraphs[0], "NorthStar Brand Copilot  ·  AI Agents on Databricks", 9, GRAY)
    notes(s, speaker)
    return s


# ----------------------------------------------------------------------------
# Slide 1 — Title
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY, rounded=False)
rect(s, 0, Inches(5.0), SW, Inches(0.12), LAVA, rounded=False)
t = textbox(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.0))
_set(t.paragraphs[0], "NorthStar Brand Copilot", 46, WHITE, bold=True)
p = t.add_paragraph(); _set(p, "An AI Agent on Databricks for Consumer Packaged Goods", 24, LAVA, bold=False)
sub = textbox(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4))
_set(sub.paragraphs[0], "LangGraph · MLflow · Genie · Vector Search · Lakebase · Databricks Apps", 16, WHITE)
p = sub.add_paragraph(); _set(p, "Powered by Claude Sonnet 4.5  ·  Field Engineering Demo", 13, RGBColor(0xB9,0xC2,0xC6))
notes(s, "Welcome. This is an end-to-end demo of authoring and deploying an AI agent on Databricks, "
         "for the CPG industry. In ~15 minutes you'll see one conversational agent answer three very "
         "different kinds of questions — documents, analytics, and remembered context — all on the "
         "Databricks platform, fully governed, traced, and evaluated. Everything you'll see is built "
         "with the latest Databricks-recommended pattern: an MLflow ResponsesAgent running inside a "
         "Databricks App, with tools attached via managed MCP servers.")

# ----------------------------------------------------------------------------
# Slide 2 — The use case / problem
# ----------------------------------------------------------------------------
content_slide(
    "The Problem", "CPG · Brand & Sales Teams",
    [("NorthStar Brands: a multi-category CPG company (snacks, beverages, personal care).", 0),
     ("Brand managers & field-sales reps need answers that span three data shapes:", 0),
     ("The numbers — sell-in/sell-out, trade-promotion ROI, inventory, market share", 1),
     ("The documents — product specs & allergens, consumer reviews, playbooks, competitive briefs", 1),
     ("Their own context — decisions made, items flagged for review", 1),
     ("Today that means bouncing between BI dashboards, a wiki, and spreadsheets.", 0),
     ("Goal: one conversational assistant that does all three — governed and observable.", 0)],
    "Set the business context. CPG teams live across three disconnected tools: BI for the numbers, "
    "a wiki/SharePoint for documents, and their own notes for context. Switching tools kills "
    "productivity and decisions get made on stale or incomplete information. The opportunity is a "
    "single agent that meets them in natural language and pulls from all three — without sacrificing "
    "governance. Emphasize this is a realistic enterprise pattern, not a toy chatbot.")

# ----------------------------------------------------------------------------
# Slide 3 — Solution overview
# ----------------------------------------------------------------------------
content_slide(
    "The Solution", "One Agent · Three Capabilities",
    [("A LangGraph agent that routes each question to the right Databricks-native tool:", 0),
     ("INSIGHTS  →  Vector Search (RAG over unstructured CPG documents)", 1),
     ("ANALYTICS →  Genie space (natural-language SQL over the data)", 1),
     ("MEMORY    →  Lakebase (remembers decisions, with semantic recall)", 1),
     ("Wrapped as an MLflow ResponsesAgent; deployed inside a Databricks App.", 0),
     ("Powered by Claude Sonnet 4.5 via Databricks Foundation Model APIs.", 0),
     ("Every request is traced (MLflow) and quality is measured (Agent Evaluation).", 0)],
    "Here's the shape of the solution. The agent is a router — it reads the question and decides which "
    "specialist tool to call: Vector Search for 'what does the document say', Genie for 'what do the "
    "numbers show', Lakebase for 'what did we decide'. The key message: these aren't bolted-on APIs — "
    "they're first-class Databricks capabilities, governed by Unity Catalog, and the agent is "
    "authored the way Databricks now recommends. Claude Sonnet 4.5 is the reasoning engine.")

# ----------------------------------------------------------------------------
# Slide 4 — Architecture (diagram)
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE, rounded=False)
rect(s, 0, 0, Inches(0.18), SH, LAVA, rounded=False)
k = textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(0.4))
_set(k.paragraphs[0], "ARCHITECTURE", 13, LAVA, bold=True)
tt = textbox(s, Inches(0.7), Inches(0.75), Inches(11.8), Inches(0.7))
_set(tt.paragraphs[0], "Agent runs inside the Databricks App", 28, NAVY, bold=True)
# App container
rect(s, Inches(1.0), Inches(1.7), Inches(11.3), Inches(1.0), LGRAY)
appx = textbox(s, Inches(1.2), Inches(1.78), Inches(11.0), Inches(0.9))
_set(appx.paragraphs[0], "Databricks App  —  React chat UI  +  MLflow AgentServer (ResponsesAgent, autolog)", 15, NAVY, bold=True)
# Agent box
ag = rect(s, Inches(4.0), Inches(2.95), Inches(5.3), Inches(0.9), NAVY)
box_text(ag, "LangGraph Agent  ·  Claude Sonnet 4.5", 16, WHITE)
# three tool boxes
ty, tw, th, gap = Inches(4.3), Inches(3.5), Inches(1.2), Inches(0.45)
xs = [Inches(1.0), Inches(1.0)+tw+gap, Inches(1.0)+2*(tw+gap)]
labels = [("Vector Search", "RAG over documents", LAVA),
          ("Genie Space", "NL → SQL analytics", BLUE),
          ("Lakebase", "Memory (Postgres)", GREEN)]
for x, (a, b, c) in zip(xs, labels):
    box = rect(s, x, ty, tw, th, c)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf.paragraphs[0], a, 16, WHITE, bold=True); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    pp = tf.add_paragraph(); _set(pp, b, 12, WHITE); pp.alignment = PP_ALIGN.CENTER
# data source labels
srcs = ["specs · reviews · playbooks · briefs", "sales · promos · inventory · share", "decisions · action items"]
for x, slabel in zip(xs, srcs):
    d = textbox(s, x, ty+th+Inches(0.15), tw, Inches(0.6))
    _set(d.paragraphs[0], slabel, 11, GRAY); d.paragraphs[0].alignment = PP_ALIGN.CENTER
# footer line
fr = textbox(s, Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.5))
_set(fr.paragraphs[0], "Tools attached via Unity-Catalog-governed MCP servers · Traces & evaluation in MLflow", 12, NAVY, bold=True)
notes(s, "Walk the architecture top-down. Everything lives in ONE deployable unit — the Databricks App — "
         "which serves the chat UI and runs the MLflow AgentServer. Inside it, the LangGraph agent uses "
         "Claude to decide routing. The three tools: Vector Search and Genie are attached as "
         "Databricks-managed MCP servers (the current recommended, governed pattern); Lakebase memory "
         "uses the managed AsyncDatabricksStore. Bottom line: the agent code, the tools, the model, and "
         "the traces all stay on Databricks under Unity Catalog governance — no data leaves the platform.")

# ----------------------------------------------------------------------------
# Slide 5 — Component map
# ----------------------------------------------------------------------------
content_slide(
    "How It Maps to Databricks", "Components",
    [("LangGraph — agent orchestration / routing", 0),
     ("MLflow ResponsesAgent + AgentServer — the agent contract & async server", 0),
     ("MLflow Autolog + Agent Evaluation — tracing and quality scoring", 0),
     ("Genie — natural-language SQL over governed tables", 0),
     ("Mosaic AI Vector Search — embeddings + retrieval for RAG", 0),
     ("Lakebase — managed Postgres for low-latency agent memory", 0),
     ("Databricks Apps + Asset Bundles — host & deploy the agent", 0),
     ("Foundation Model APIs — Claude Sonnet 4.5 + GTE embeddings", 0),
     ("Unity Catalog — governs every table, index, function, and grant", 0)],
    "This slide is the 'checklist' — every capability the customer asked about, mapped to a concrete "
    "Databricks product. The point to land: these compose natively. You don't stitch together five "
    "vendors; the orchestration framework (LangGraph) plugs into Databricks tools, MLflow gives you "
    "the production contract and observability, and Unity Catalog governs all of it. Pause here if the "
    "audience is technical — invite questions on any specific component before the deep dives.")

# ----------------------------------------------------------------------------
# Slide 6 — Data foundation
# ----------------------------------------------------------------------------
content_slide(
    "Data Foundation", "Phase 1 · Synthetic CPG Data",
    [("8 Delta tables in Unity Catalog (serverless_stable…northstar_cpg):", 0),
     ("Structured: products (52 SKUs), retailers (10), sales_facts (~41k weekly rows),", 1),
     ("inventory, distribution, trade_promotions (600, with ROI), market_share", 1),
     ("Unstructured: documents (83) — authored specs, reviews, playbooks, briefs", 1),
     ("Realistic & internally consistent: e.g. BOGO/TPR promos skew negative-ROI,", 0),
     ("matching the trade-promotion playbook the agent later cites.", 1),
     ("Change Data Feed enabled on documents for Vector Search sync.", 0)],
    "Quickly establish credibility of the data. It's synthetic but deliberately coherent — the numbers "
    "and the documents agree, so cross-tool questions land. For example, the trade promotions are "
    "modeled so that BOGO and deep-discount TPRs tend to lose money, which is exactly what the "
    "authored playbook document recommends avoiding. That consistency is what makes the multi-tool "
    "demo (Genie finds the losers, Vector Search explains the fix) feel real. All data is Delta in "
    "Unity Catalog, generated on serverless compute.")

# ----------------------------------------------------------------------------
# Slide 7 — Vector Search
# ----------------------------------------------------------------------------
content_slide(
    "INSIGHTS — Vector Search", "RAG over unstructured documents",
    [("Delta-synced Vector Search index over the documents table.", 0),
     ("Whole-document embeddings via managed databricks-gte-large-en (1024-dim).", 0),
     ("Metadata (doc_type, title, brand, category) synced for filtering & citations.", 0),
     ("Demo question: “What allergens are in Aurora Oat Milk, and what do consumers dislike?”", 0),
     ("Returns: oats / may contain tree nuts + ‘separates in coffee’, ‘thin texture’ — cited.", 1),
     ("Grounded answers with sources — not hallucination.", 0)],
    "The Insights capability. Vector Search indexes the unstructured documents and keeps them in sync "
    "with the Delta table automatically. We use Databricks-managed embeddings so there are no vectors "
    "to maintain. The demo question pulls allergen facts from a spec sheet AND complaints from review "
    "text, then cites both. Stress 'grounded + cited' — that's what makes RAG trustworthy for "
    "enterprise use. Mention metadata is retained so you can filter (e.g. only reviews for one brand) "
    "and show provenance.")

# ----------------------------------------------------------------------------
# Slide 8 — Genie
# ----------------------------------------------------------------------------
content_slide(
    "ANALYTICS — Genie", "Natural-language SQL over the numbers",
    [("A Genie space scoped to the 7 structured tables, with curated instructions.", 0),
     ("Knows business semantics: ‘last quarter’ = 13 weeks, sell-through, ROI, joins.", 1),
     ("Demo question: “Which promotions had the most negative ROI last quarter?”", 0),
     ("Genie writes & runs the SQL, returns a ranked table + a takeaway.", 1),
     ("Attached to the agent as a managed MCP server — no SQL written by hand.", 0),
     ("Combo: pair with Vector Search to explain WHY (cite the playbook).", 0)],
    "The Analytics capability. Same chat box, but a quantitative question routes to a Genie space that "
    "generates and executes SQL across the sales and promotion tables. The reason it's reliable is the "
    "curated instructions and example SQL we gave the space — it understands NorthStar's definitions "
    "of 'last quarter', sell-through, and ROI, and the correct joins. Show the optional combo: ask it "
    "to also consult the playbook, and the agent calls Genie AND Vector Search in one turn — connecting "
    "the data to the recommended action. That's the multi-agent payoff.")

# ----------------------------------------------------------------------------
# Slide 9 — Lakebase
# ----------------------------------------------------------------------------
content_slide(
    "MEMORY — Lakebase", "The agent remembers",
    [("Lakebase = Databricks managed Postgres (OLTP) for agent memory.", 0),
     ("Long-term memory via AsyncDatabricksStore with semantic recall.", 0),
     ("Demo: “Remember we’re cutting BOGO at Walgreens…”  →  saved.", 0),
     ("Later, in a new message: “What did we decide about the West region?” → recalled.", 0),
     ("Low-latency, transactional state that lives next to the lakehouse.", 0),
     ("Connection & token rotation handled by the managed store.", 1)],
    "The Memory capability — the one most chatbots lack. Decisions and action items persist to "
    "Lakebase, Databricks' managed Postgres, and are recalled semantically using the same embedding "
    "model behind Vector Search. Demo it live in two turns: save a decision, then ask about it in a "
    "fresh message and watch it come back. The architectural point: this is OLTP — fast, "
    "transactional, per-user state — sitting right next to your analytical lakehouse, which is exactly "
    "what agentic apps need and historically required a separate database to get.")

# ----------------------------------------------------------------------------
# Slide 9a — Lakebase deep dive: what it is
# ----------------------------------------------------------------------------
content_slide(
    "Lakebase — Postgres for the Lakehouse", "Deep Dive · Why it matters",
    [("A fully managed, Postgres-compatible OLTP database, native to Databricks.", 0),
     ("Built on serverless Postgres — compute & storage are separated for elasticity.", 0),
     ("Brings transactional, low-latency data right next to your analytical lakehouse:", 0),
     ("no separate database vendor, no brittle reverse-ETL, one security model", 1),
     ("Purpose-built for operational apps, ML feature serving, and AI-agent memory/state.", 0),
     ("Governed by Unity Catalog; authenticated with Databricks identity (OAuth) —", 0),
     ("no static credentials or long-lived secrets to manage.", 1)],
    "Frame the 'why' before the features. Historically, putting an app or an agent on top of your "
    "lakehouse meant bolting on a separate operational database — RDS/Postgres/Dynamo — with "
    "reverse-ETL pipelines to move data out, a second security model, and separate ops. Lakebase "
    "collapses that: it's managed Postgres that lives on the Databricks platform, governed by Unity "
    "Catalog, accessed with your Databricks identity. For agentic apps this is the missing piece — the "
    "agent's transactional state and the analytical data it reasons over finally live on one governed "
    "platform. It's standard Postgres, so existing skills, drivers, and tools carry over.")

# ----------------------------------------------------------------------------
# Slide 9b — Lakebase key capabilities
# ----------------------------------------------------------------------------
content_slide(
    "Lakebase — Key Capabilities", "Deep Dive · Features",
    [("Serverless & elastic — autoscaling compute, scales to zero, pay-per-use.", 0),
     ("Low latency & high concurrency — single-digit-millisecond reads, thousands of QPS.", 0),
     ("Instant branching & copy-on-write clones — spin up dev/test copies in seconds.", 0),
     ("Lakehouse sync (online tables) — keep Delta/UC tables synced to Postgres for serving.", 0),
     ("Full Postgres ecosystem — standard drivers/tools, extensions incl. pgvector.", 0),
     ("Enterprise security — Unity Catalog governance, rotating OAuth tokens, encryption.", 0),
     ("Built for AI — agent memory, feature serving, and app backends on Databricks Apps.", 0)],
    "Now the capabilities — emphasize the ones that matter for agents. Low latency: the assistant "
    "recalls a saved decision in milliseconds. Branching: copy-on-write storage lets you clone "
    "production data instantly for safe dev/test — a superpower you don't get from a bolted-on RDS. "
    "Lakehouse sync / online tables: serve governed lakehouse data to the app without building "
    "pipelines. pgvector: semantic memory and retrieval live right in Postgres. OAuth: the app's "
    "service principal authenticates with short-lived, rotating tokens — no secrets in code. Tie it "
    "back to this demo: the Assistant tab's memory runs on Lakebase via the managed "
    "AsyncDatabricksStore — save a decision, recall it later by meaning, all under the app's Databricks "
    "identity. That's an enterprise-grade operational layer for agents, native to the platform.")

# ----------------------------------------------------------------------------
# Slide 10 — The agent
# ----------------------------------------------------------------------------
content_slide(
    "The Agent", "LangGraph + MLflow ResponsesAgent",
    [("LangGraph create_agent with Claude Sonnet 4.5 as the router/reasoner.", 0),
     ("Tools: Genie (MCP) + Vector Search (MCP) + Lakebase memory tools.", 0),
     ("Wrapped in MLflow ResponsesAgent — instant compatibility with Playground,", 0),
     ("Agent Evaluation, and Apps deployment; auto-inferred model signature.", 1),
     ("mlflow.langchain.autolog() captures the full trace of every request.", 0),
     ("Graceful degradation + clear tool descriptions + LLM timeouts (best practice).", 0)],
    "Now the agent itself. It's a LangGraph agent — the supervisor/router is the LLM choosing tools. "
    "We wrap it in MLflow's ResponsesAgent, which is the interface Databricks recommends: it gives you "
    "out-of-the-box compatibility with the AI Playground, Agent Evaluation, and App deployment, and "
    "MLflow infers the model signature for you. Autolog gives full tracing for free. Mention the "
    "engineering best practices we followed: clear tool descriptions so routing is accurate, graceful "
    "fallback if a tool is unavailable, and timeouts on LLM calls.")

# ----------------------------------------------------------------------------
# Slide 11 — Deployment
# ----------------------------------------------------------------------------
content_slide(
    "Deployment — Databricks Apps", "The latest recommended pattern",
    [("Agent runs INSIDE a Databricks App (ResponsesAgent + AgentServer).", 0),
     ("Deployed via Databricks Asset Bundles (databricks bundle deploy / run).", 0),
     ("NOT the older log-to-UC → Model Serving endpoint pattern.", 1),
     ("Why: rapid local iteration, git/CI-CD, async scaling, built-in chat UI.", 0),
     ("Resource grants in databricks.yml: endpoints, Genie, index, Lakebase.", 0),
     ("One-command deploy + service-principal grants in the deployment/ folder.", 0)],
    "Deployment is where we follow the May-2026 Databricks guidance exactly: for new agents, run them "
    "inside a Databricks App via Asset Bundles, not as a Model Serving endpoint. The advantages are "
    "rapid local iteration, git-based versioning and CI/CD, async scalability, and a built-in chat UI. "
    "All the least-privilege grants for the app's service principal are declared in databricks.yml and "
    "applied by a single deploy script. If asked, note this consolidates the 'agent' and the "
    "'frontend' into one governed, versioned deployable.")

# ----------------------------------------------------------------------------
# Slide 12 — Observability
# ----------------------------------------------------------------------------
content_slide(
    "Observability — MLflow Tracing", "See how the agent decided",
    [("Every request is a trace: supervisor → tool → LLM, with timings & tokens.", 0),
     ("Genie route:  agent → ChatDatabricks → query_space… [Genie] → ChatDatabricks", 1),
     ("Vector Search route:  agent → ChatDatabricks → documents_index [VS] → …", 1),
     ("Same agent, different routes — visible side-by-side in the MLflow UI.", 0),
     ("Powered by mlflow.langchain.autolog() — no manual instrumentation.", 0),
     ("(Trace data is stored & queryable for debugging and monitoring.)", 0)],
    "This is the 'trace reveal' — open MLflow and show two traces side by side. For a numbers question "
    "you'll see the Genie tool fire; for a document question you'll see the Vector Search tool. You can "
    "drill into any span to see the exact SQL, the retrieved documents, latency, and token usage. This "
    "is how you debug and monitor an agent in production. It all comes from one line — autolog. "
    "Presenter tip: generate fresh traces from the provided notebook beforehand so they're ready.")

# ----------------------------------------------------------------------------
# Slide 13 — Evaluation
# ----------------------------------------------------------------------------
content_slide(
    "Quality — MLflow Agent Evaluation", "Measure before you ship",
    [("LLM-judge scorers run against a curated CPG question set (with expected facts).", 0),
     ("Results on the 10-question set:", 0),
     ("Safety 1.0   ·   Relevance 0.9   ·   Correctness 0.8", 1),
     ("Per-question judge rationale attached to each trace.", 0),
     ("Run it in CI to gate changes before redeploying the app.", 0)],
    "Quality. We don't ship and hope — MLflow Agent Evaluation runs LLM-judge scorers (correctness vs. "
    "expected facts, relevance, safety) over a curated question set and attaches the judge's reasoning "
    "to each trace. Our scores: Safety 1.0, Relevance 0.9, Correctness 0.8 — strong, with clear places "
    "to improve. The workflow point: wire this into CI so a regression in answer quality blocks the "
    "deploy, just like a failing unit test. One gotcha worth knowing: the built-in judges require the "
    "databricks-agents package to run.")

# ----------------------------------------------------------------------------
# Slide 14 — Demo flow
# ----------------------------------------------------------------------------
content_slide(
    "Live Demo Flow", "Three acts + the reveal",
    [("Act 1 — INSIGHTS: allergens & dislikes for Aurora Oat Milk (Vector Search).", 0),
     ("Act 2 — ANALYTICS: worst-ROI promotions last quarter (Genie).", 0),
     ("Act 3 — MEMORY: save a West-region decision, then recall it (Lakebase).", 0),
     ("Reveal — open MLflow: routing traces + evaluation scores.", 0),
     ("Close — one governed agent, three tools, all on Databricks.", 0)],
    "This is your run-of-show. Keep it tight: one question per capability, then the MLflow reveal. The "
    "narrative arc is documents → numbers → memory → 'and here's how we know it works and how it "
    "decided'. If you have time, do the combo question in Act 2 to show multi-tool reasoning. The full "
    "scripted version with exact questions and expected answers is in DEMO_SCRIPT.md.")

# ----------------------------------------------------------------------------
# Slide 15 — Why Databricks / close
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY, rounded=False)
rect(s, 0, Inches(1.7), Inches(2.0), Inches(0.10), LAVA, rounded=False)
t = textbox(s, Inches(0.9), Inches(0.8), Inches(11.5), Inches(0.9))
_set(t.paragraphs[0], "Why Databricks", 34, WHITE, bold=True)
body = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
pts = ["One platform: data, embeddings, NL-SQL, OLTP, model, agent, and observability.",
       "Governed by Unity Catalog end-to-end — no data leaves the platform.",
       "Authored & deployed the Databricks-recommended way (ResponsesAgent on Apps).",
       "Traced and evaluated with MLflow — production-ready, not a prototype.",
       "Fully reproducible from the repo: data → tools → agent → eval → one-command deploy."]
first = True
for x in pts:
    p = body.paragraphs[0] if first else body.add_paragraph(); first = False
    _set(p, "•  " + x, 18, WHITE); p.space_after = Pt(12)
    p.runs[0].font.color.rgb = WHITE
end = textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6))
_set(end.paragraphs[0], "Thank you  ·  Questions?", 18, LAVA, bold=True)
notes(s, "Close on the platform story. The reason this demo is compelling isn't any single feature — "
         "it's that data, retrieval, NL-SQL, transactional memory, the model, the agent, and the "
         "observability all live on one governed platform. That collapses the integration burden and "
         "keeps data in place under Unity Catalog. Reiterate it's production-shaped — traced, "
         "evaluated, bundle-deployed — and fully reproducible. Then open for questions.")

# COMMAND ----------
import shutil
os.makedirs(OUT_DIR, exist_ok=True)
# Save to local disk first (UC Volume FUSE mounts reject pptx's random-access writes),
# then stream-copy the finished file onto the Volume.
TMP = "/tmp/NorthStar_Brand_Copilot_Demo.pptx"
prs.save(TMP)
shutil.copyfile(TMP, OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst), "| bytes:", os.path.getsize(OUT))
dbutils.notebook.exit(OUT)
